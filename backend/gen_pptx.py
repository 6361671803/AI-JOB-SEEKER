from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

SCRATCH = Path(r"C:\Users\kmdri\AppData\Local\Temp\claude\c--job-seeker-agent\25df0348-e594-4ace-988c-2c1eb7254241\scratchpad")
OUT = Path(r"C:\Users\kmdri\Downloads\Agentic_Job_Finder_Presentation.pptx")

NAVY = RGBColor(0x10, 0x16, 0x2A)
INDIGO = RGBColor(0x45, 0x50, 0xE5)
MUTED = RGBColor(0x5B, 0x64, 0x78)
LIGHT_BG = RGBColor(0xEE, 0xF0, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1C, 0x7A, 0x4D)
AMBER = RGBColor(0x92, 0x65, 0x0A)
RED = RGBColor(0xB8, 0x28, 0x3A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_title(slide, text, subtitle=None, dark=False):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE if dark else NAVY
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xC5, 0xC9, 0xE0) if dark else MUTED
    # accent underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.15), Inches(2.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = INDIGO
    line.line.fill.background()
    line.shadow.inherit = False


def bullets(slide, items, left=Inches(0.7), top=Inches(1.5), width=Inches(11.9), height=Inches(5.4), size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(size - level * 2)
        p.font.color.rgb = NAVY if level == 0 else MUTED
        p.space_after = Pt(8)


def add_image_slide(title_text, img_path, caption=None):
    s = add_slide()
    add_bg(s)
    add_title(s, title_text)
    pic_path = str(SCRATCH / img_path)
    pic = s.shapes.add_picture(pic_path, Inches(0.8), Inches(1.5), width=Inches(11.7))
    if pic.height > Inches(5.4):
        # rescale by height if too tall
        ratio = Inches(5.4) / pic.height
        pic.height = Inches(5.4)
        pic.width = Emu(int(pic.width * ratio))
        pic.left = Inches((13.333 - pic.width / 914400) / 2 * 914400) if False else pic.left
    if caption:
        box = s.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = MUTED
        p.alignment = PP_ALIGN.CENTER
    return s


def add_table_slide(title_text, headers, rows, col_widths=None):
    s = add_slide()
    add_bg(s)
    add_title(s, title_text)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    left, top, width, height = Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.3)
    table_shape = s.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(13)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = INDIGO
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(11.5)
            cell.text_frame.paragraphs[0].font.color.rgb = NAVY
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 0 else LIGHT_BG
    return s


# ============ SLIDE 1: TITLE ============
s = add_slide()
add_bg(s, NAVY)
badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(1.6), Inches(1.5), Inches(1.5))
badge.fill.solid(); badge.fill.fore_color.rgb = INDIGO; badge.line.fill.background()
tf = badge.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "AI"; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

box = s.shapes.add_textbox(Inches(1), Inches(3.3), Inches(11.3), Inches(1.6))
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Agentic Job Finder and Application Assistant"
p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

box2 = s.shapes.add_textbox(Inches(1.5), Inches(4.6), Inches(10.3), Inches(1.0))
tf2 = box2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "A Multi-Agent AI System for Automated Job Discovery, Resume Matching, and Application Preparation"
p2.font.size = Pt(16); p2.font.color.rgb = RGBColor(0xC5, 0xC9, 0xE0)
p2.alignment = PP_ALIGN.CENTER

box3 = s.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.5))
p3 = box3.text_frame.paragraphs[0]
p3.text = "Python · FastAPI · React · Ollama · Playwright · Tavily"
p3.font.size = Pt(13); p3.font.color.rgb = RGBColor(0x8A, 0x91, 0xC7)
p3.alignment = PP_ALIGN.CENTER

# ============ SLIDE 2: OBJECTIVES ============
s = add_slide(); add_bg(s); add_title(s, "Objectives", "What the system is built to do")
bullets(s, [
    ("Extract structured candidate data from a resume without inventing information", 0),
    ("Discover real companies and their official career pages via live web search", 0),
    ("Discover real, current job openings by rendering each career page with a real browser", 0),
    ("Filter jobs by a configurable posting-recency window (default 30 days)", 0),
    ("Score every job against the resume with an explainable, evidence-based match %", 0),
    ("Require two independent human approvals before any application work begins or submits", 0),
    ("Automate application form-filling — stop automatically on CAPTCHA / OTP / login", 0),
    ("Never click a final Submit button — the user always does that personally", 0),
    ("Track every job's status end-to-end in a persistent application tracker", 0),
])

# ============ SLIDE 3: TECH STACK ============
add_table_slide("Technology Stack", ["Layer", "Technology", "Purpose"], [
    ["Backend", "Python 3.10, FastAPI, Uvicorn", "REST API server"],
    ["Database", "SQLite + SQLAlchemy ORM", "Persistence"],
    ["Validation", "Pydantic / pydantic-settings", "Schemas & config"],
    ["Resume Parsing", "pypdf, python-docx", "PDF / DOCX text extraction"],
    ["LLM", "Gemini (cloud, default) / Ollama (local) / OpenAI — one swappable client", "Structured extraction & reasoning"],
    ["Web Search", "Tavily Search API", "Company & job discovery"],
    ["Browser Automation", "Playwright (Chromium)", "Render pages, fill forms, screenshot"],
    ["Frontend", "React 18 (Vite)", "Single-page application UI"],
    ["Styling", "Custom CSS design system", "Consistent visual language"],
], col_widths=[3.0, 4.4, 4.5])

# ============ SLIDE 4: ALGORITHMS ============
s = add_slide(); add_bg(s); add_title(s, "Algorithms & Techniques", "Engineering choices that prevent AI hallucination")
bullets(s, [
    ("JSON-Schema-Constrained LLM Extraction", 0),
    ("Every LLM call forces strict, schema-valid JSON output — never free-form prose", 1),
    ("Deterministic Whole-Word Skill Matching", 0),
    ("Regex word-boundary matching — zero hallucination risk on \"do I have this skill\"", 1),
    ("Weighted Multi-Criteria Match Scoring", 0),
    ("35% skills + 20% experience + 15% education + 15% role + 15% location", 1),
    ("Grounding / Anti-Hallucination Backstops", 0),
    ("URLs, dates, and locations must be verified against real scraped page content", 1),
    ("Two-Stage Job Extraction & Index-Based Link Selection", 0),
    ("Prevents the LLM from fabricating data that isn't on the page it's reading", 1),
    ("Finite State Machine Workflow", 0),
    ("Candidate.state / Job.status drive orchestration and gate human approvals", 1),
], size=15)

# ============ SLIDE 5: ARCHITECTURE (image) ============
add_image_slide("System Architecture", "svg_architecture.png",
                 "All external calls (LLM, search, browser) happen server-side in FastAPI.")

# ============ SLIDE 6: WORKFLOW PIPELINE (image) ============
add_image_slide("11-Phase Agent Pipeline", "svg_pipeline.png",
                 "DISCOVER -> MATCH -> RANK -> USER SELECTS -> PREPARE -> STOP -> USER REVIEWS -> USER SUBMITS")

# ============ SLIDE 7: ER DIAGRAM (image) ============
add_image_slide("Database Design (ER Diagram)", "svg_er_diagram.png",
                 "Candidate 1:N Company, Candidate 1:N Job, Company 1:N Job")

# ============ SLIDE 8: STATE MACHINE (image) ============
add_image_slide("Job Status State Machine", "svg_state_machine.png",
                 "Two starred transitions require an explicit, independent human approval")

# ============ SLIDE 9: BACKEND FILES ============
add_table_slide("Backend — Key Files", ["File", "Purpose"], [
    ["app/main.py", "All FastAPI REST endpoints"],
    ["app/db/models.py", "Candidate, Company, Job ORM models"],
    ["app/services/llm_client.py", "LLM provider abstraction, all prompts & schemas"],
    ["app/services/browser_client.py", "Playwright page render & form-filling"],
    ["app/services/skill_matcher.py", "Deterministic skill matching"],
    ["app/agents/resume_analyzer.py", "Resume -> structured profile"],
    ["app/agents/company_discovery_agent.py", "Company + careers page discovery"],
    ["app/agents/job_discovery_agent.py", "Two-stage job listing/detail extraction"],
    ["app/agents/matching_agent.py", "Resume/job match scoring"],
    ["app/agents/application_preparation_agent.py", "Form-fill automation, never submits"],
], col_widths=[4.8, 7.1])

# ============ SLIDE 10: FRONTEND FILES ============
add_table_slide("Frontend — Key Files", ["File", "Purpose"], [
    ["src/App.jsx", "Top-level view/state machine & navigation"],
    ["src/api/client.js", "Every backend API call (typed fetch wrappers)"],
    ["components/pages/Dashboard.jsx", "Live stats & AI search status"],
    ["components/JobResults.jsx", "Ranked, filterable, sortable job dashboard"],
    ["components/shared/MatchDetailsModal.jsx", "Full match breakdown view"],
    ["components/shared/JobDescriptionModal.jsx", "Full job description view"],
    ["components/ApplicationPreparation.jsx", "Final review & Approval #2"],
    ["components/ApplicationTracker.jsx", "Full lifecycle tracking table"],
], col_widths=[4.8, 7.1])

# ============ SLIDE 11: HOW TO RUN ============
s = add_slide(); add_bg(s); add_title(s, "How to Run the Project")
box = s.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.9), Inches(5))
tf = box.text_frame; tf.word_wrap = True
lines = [
    "Prerequisites:",
    "  Python 3.10+, Node.js 18+, Ollama (running, model pulled), Tavily API key",
    "",
    "One-command startup (from project root):",
    "  npm run dev",
    "",
    "This starts both servers together via `concurrently`:",
    "  Backend:  uvicorn app.main:app --port 8000",
    "  Frontend: vite dev server on port 5173",
    "",
    "Configuration lives in backend/.env:",
    "  LLM_PROVIDER=ollama",
    "  LLM_MODEL=qwen2.5:3b",
    "  TAVILY_API_KEY=<your key>",
]
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    is_cmd = line.strip().startswith(("npm", "uvicorn", "Backend:", "Frontend:", "LLM_", "TAVILY_"))
    p.font.name = "Consolas" if is_cmd else "Segoe UI"
    p.font.size = Pt(14)
    p.font.color.rgb = INDIGO if is_cmd else NAVY
    p.font.bold = line.endswith(":") and not is_cmd

# ============ SLIDE 12: ERRORS & FIXES ============
add_table_slide("Key Errors Encountered & Fixes", ["Problem", "Root Cause -> Fix"], [
    ["openai/httpx crash on startup", "Version mismatch -> pinned httpx==0.27.2"],
    ["OpenAI quota exhausted", "No billing configured -> switched to local Ollama"],
    ["Resume hallucination (fake degree)", "Weak prompt -> explicit \"never invent\" rules, verified fix"],
    ["Job data fabricated (fake dates, repeated jobs)", "LLM asked for data not on that page -> two-stage extraction"],
    ["Fabricated application URLs", "Free-text URL generation -> index-based link selection"],
    ["\"0 companies found\" silently", "Expired Tavily key (403) -> surfaced as explicit error"],
    ["Stale dev server serving old code", "uvicorn --reload leaked processes -> removed --reload"],
    ["Gemini model deprecated (404)", "Queried live model list instead of guessing -> found replacement"],
    ["Gemini free tier only 20 req/day", "Switched to gemini-flash-lite-latest -> verified full run, 0 errors"],
], col_widths=[5.4, 6.5])

# ============ SLIDE 13: POST-LAUNCH ENHANCEMENTS ============
s = add_slide(); add_bg(s); add_title(s, "Post-Launch Enhancements", "User-driven improvements after the initial 11-phase build")
bullets(s, [
    ("Gemini cloud LLM provider added (alongside Ollama & OpenAI)", 0),
    ("Full resume extraction: 215s (local) -> 19.8s (Gemini) — ~11x faster, verified on identical input", 1),
    ("Same hallucination tests from Section 12 re-run and confirmed clean on Gemini", 1),
    ("Wider discovery coverage", 0),
    ("Company cap raised 15 -> 25; found & fixed a hidden 10-name-per-query extraction ceiling", 1),
    ("Job-scraping was silently limited to the first 8 companies -> raised to all 25", 1),
    ("Official website / careers page accuracy", 0),
    ("Retry search + domain-derivation from verified URLs -> every company now resolves a website", 1),
], size=16)

# ============ SLIDE 14: TERMINAL LOG HIGHLIGHTS ============
add_table_slide("Terminal Log Highlights (Real Input & Output)", ["Command", "Result"], [
    ["curl .../resume/upload (before fix)", "Internal Server Error -> httpx 'proxies' TypeError"],
    ["web_search('...AI Engineer...')", "403 Forbidden -> expired Tavily key found & replaced"],
    ["curl .../discover-jobs (before fix)", "Same fake job repeated 10x, fake 2023 date"],
    ["curl .../discover-jobs (after fix)", "Real distinct jobs, verified real URLs"],
    ["curl .../match-jobs", "Strong match 86% vs Weak match 17% — correctly discriminated"],
    ["curl .../prepare-application (CAPTCHA)", "Stopped, 0 fields touched — screenshot confirmed"],
    ["client.models.list() on Gemini", "50+ real models queried live, picked flash-lite-latest"],
    ["curl .../discover-companies (final)", "HTTP 200 — 25 companies, 0 missing website"],
], col_widths=[5.6, 6.3])

# ============ SLIDE 15: CONCLUSION ============
s = add_slide(); add_bg(s); add_title(s, "Conclusion & Future Enhancements")
bullets(s, [
    ("All 11 phases implemented end-to-end, verified against live data throughout", 0),
    ("Key lesson: small local LLMs need deterministic, code-level grounding — not just better prompts", 0),
    ("Human-in-the-loop approval gates preserved at every consequential step", 0),
    ("Post-launch: cloud LLM (Gemini) solved speed, but introduced its own failure mode (quota) — verified live", 0),
    ("Future enhancements:", 0),
    ("Streaming/incremental progress via SSE or WebSockets", 1),
    ("Conversational AI assistant panel over discovered job data", 1),
    ("Persistent cross-candidate company-URL cache to reduce repeat resolution calls", 1),
    ("Dark-mode theme toggle using the existing design-token system", 1),
    ("Automated end-to-end test suite over recorded fixture pages", 1),
], size=16)

prs.save(str(OUT))
print("Saved:", OUT, OUT.stat().st_size, "bytes")
